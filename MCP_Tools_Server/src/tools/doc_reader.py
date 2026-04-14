import asyncio
import os
import shutil
import time
import uuid
from pathlib import Path
from typing import Dict, Tuple, Union
from urllib.parse import urlparse

import aiofiles
import aiohttp
from docx import Document
from ollama import AsyncClient
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from sse.event_bus import event_bus
from utils.logger.AgentLogger import quickLog
from utils.task_scheduler import scheduler

TEMP_DIR = Path(__file__).parent.parent / "temp"
TEMP_DIR.mkdir(parents=True, exist_ok=True)


# -------------------- HELPERS --------------------


def estimate_tokens(text: str) -> int:
    return max(1, len(text) // 4)


def first_100_words(text: str) -> str:
    return " ".join(text.split()[:100])


def clean_text(text: str) -> str:
    lines = [line.strip() for line in text.splitlines()]
    seen = set()
    cleaned = []

    for line in lines:
        if not line or len(line) < 3:
            continue
        if len(line.split()) < 2:
            continue

        key = line.lower()
        if key in seen:
            continue

        seen.add(key)
        cleaned.append(line)

    return "\n".join(cleaned)


def ordinal(n: int) -> str:
    if 10 <= n % 100 <= 20:
        suffix = "th"
    else:
        suffix = {1: "st", 2: "nd", 3: "rd"}.get(n % 10, "th")
    return f"{n}{suffix}"


def generate_readable_name(index: int, ext: str, content: str) -> str:
    words = content.split()[:10]

    safe_words = [
        w.lower().strip(".,!?()[]{}\"'").replace("/", "") for w in words if w.strip()
    ]

    joined = "_".join(safe_words) or "empty"

    return f"{ordinal(index)}_file_{ext}_{joined}.{ext}"


def generate_temp_path(ext: str) -> Path:
    return TEMP_DIR / f"{uuid.uuid4()}.{ext}"


# -------------------- FILE HANDLING --------------------


async def download_to_temp(url: str, forced_ext: str | None = None) -> Path:
    async with aiohttp.ClientSession() as session:
        async with session.get(url) as resp:
            resp.raise_for_status()

            content_type = (resp.headers.get("content-type", "") or "").lower()

            ext = forced_ext
            if not ext:
                parsed_url = urlparse(url)
                suffix_ext = Path(parsed_url.path).suffix.lstrip(".").lower()

                # Prefer explicit extension from the URL path when present.
                if suffix_ext in {
                    "pdf",
                    "docx",
                    "pptx",
                    "xlsx",
                    "txt",
                    "md",
                    "csv",
                    "json",
                    "xml",
                    "log",
                }:
                    ext = suffix_ext
                elif "pdf" in content_type:
                    ext = "pdf"
                elif "word" in content_type:
                    ext = "docx"
                elif "presentation" in content_type:
                    ext = "pptx"
                elif "spreadsheet" in content_type:
                    ext = "xlsx"
                elif "markdown" in content_type:
                    ext = "md"
                elif "text/" in content_type:
                    ext = "txt"
                elif "json" in content_type:
                    ext = "json"
                elif "xml" in content_type:
                    ext = "xml"
                elif "csv" in content_type:
                    ext = "csv"
                else:
                    ext = "bin"

            temp_path = generate_temp_path(ext)

            async with aiofiles.open(temp_path, "wb") as f:
                await f.write(await resp.read())

    print(f"[DOWNLOADED] -> {temp_path}")
    return temp_path


async def copy_to_temp(local_path: Path) -> Path:
    ext = local_path.suffix.replace(".", "")
    temp_path = generate_temp_path(ext)

    loop = asyncio.get_event_loop()
    await loop.run_in_executor(None, shutil.copy, local_path, temp_path)

    print(f"[COPIED] -> {temp_path}")
    return temp_path


async def prepare_file(item: Union[str, Tuple[str, str]]) -> Tuple[Path | None, bool, str]:
    forced_type: str | None = None

    if isinstance(item, tuple):
        item, forced_type = item

    if isinstance(item, str) and item.startswith("http"):
        path = await download_to_temp(item, forced_type)
        return path, True, item
    else:
        path = Path(item)
        if not path.exists():
            return None, False, item

        temp_path = await copy_to_temp(path)
        return temp_path, False, path.name


# -------------------- EXTRACTION --------------------


async def extract_pdf(path):
    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


async def extract_docx(path):
    doc = Document(path)
    content = []

    for p in doc.paragraphs:
        if p.text.strip():
            content.append(p.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                content.append(row_text)

    return "\n".join(content)


async def extract_pptx(path):
    prs = Presentation(path)
    slides = []

    for i, slide in enumerate(prs.slides):
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

        slides.append(f"[Slide {i + 1}]\n" + "\n".join(text))

    return "\n\n".join(slides)


async def extract_xlsx(path):
    wb = load_workbook(path, data_only=True)
    sheets = []

    for sheet in wb:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            rows.append(" | ".join(str(c) if c else "" for c in row))

        sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))

    return "\n\n".join(sheets)


async def extract_txt_md(path):
    async with aiofiles.open(path, "r", encoding="utf-8", errors="ignore") as f:
        return await f.read()


async def extract_text(path: Path):
    p = str(path).lower()

    if p.endswith(".pdf"):
        return await extract_pdf(path)
    elif p.endswith(".docx"):
        return await extract_docx(path)
    elif p.endswith(".pptx"):
        return await extract_pptx(path)
    elif p.endswith(".xlsx"):
        return await extract_xlsx(path)
    elif (
        p.endswith(".txt")
        or p.endswith(".md")
        or p.endswith(".csv")
        or p.endswith(".json")
        or p.endswith(".xml")
        or p.endswith(".log")
    ):
        return await extract_txt_md(path)
    else:
        return None


# -------------------- OLLAMA --------------------

SYSTEM_PROMPT = """You are an intelligent document summarizer. Analyze content and produce a concise summary. Prefer keeping response under 100 words but prioritize clarity."""


async def summarize(text: str, host: str):
    client = AsyncClient(host=host)

    short_text = first_100_words(text)

    start = time.time()

    response = await client.generate(
        model="gemma4:e2b",
        prompt=short_text,
        system=SYSTEM_PROMPT,
        think=False,
        keep_alive=True,
    )

    end = time.time()

    output = response.get("response", "").strip()

    return {
        "summary": output,
        "tokens": estimate_tokens(short_text) + estimate_tokens(output),
        "time": round(end - start, 3),
    }


# -------------------- MAIN --------------------


async def process_files(inputs, ollama_host: str) -> Dict:
    total_files = len(inputs)
    success_count = 0
    results = {}

    total_tokens = 0
    total_time = 0

    async def process_single(item, index):
        nonlocal success_count, total_tokens, total_time

        try:
            temp_path, is_url, original_name = await prepare_file(item)

            if temp_path is None:
                results[str(item)] = "fail to process this file"
                return

            content = await extract_text(temp_path)

            if content is None:
                results[str(item)] = "fail to process this file"
                return

            content = clean_text(content)

            if not content or len(content.strip()) < 20:
                results[str(item)] = "fail to process this file"
                return

            ext = temp_path.suffix.replace(".", "")

            if is_url:
                filename = generate_readable_name(index, ext, content)
            else:
                filename = original_name

            res = await summarize(content, ollama_host)

            results[filename] = res["summary"]
            success_count += 1
            total_tokens += res["tokens"]
            total_time += res["time"]

        except Exception:
            results[str(item)] = "fail to process this file"

        await event_bus.broadcast(
            message={
                "msg": "Processed file",
                "tool_param": str(item),
                "tool_result": results.get(str(item)),
            }
        )

    await asyncio.gather(*(process_single(i, idx + 1) for idx, i in enumerate(inputs)))

    return {
        "total_files": total_files,
        "succeed": success_count,
        "total_tokens_used": total_tokens,
        "total_time_taken": round(total_time, 3),
        "content": results,
    }


# -------------------- DUMMY USAGE --------------------

if __name__ == "__main__":
    files = [
        ("https://arxiv.org/pdf/2604.02315", "pdf"),
        "/home/pixelthreader/Desktop/pixelThreader's Workspace/OpenSource Projects/Deep-Researcher-V2/MCP_Tools_Server/src/tools/test.md",
        "/home/pixelthreader/Desktop/pixelThreader's Workspace/OpenSource Projects/Deep-Researcher-V2/MCP_Tools_Server/src/tools/test.txt",
        "/home/pixelthreader/Downloads/DOCX_TestPage.docx",
        "/home/pixelthreader/Downloads/demo.docx",
    ]

    async def main():
        result = await process_files(files, "http://localhost:11434")
        print(result)

    asyncio.run(main())
